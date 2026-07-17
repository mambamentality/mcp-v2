# tools/extract_backup_content.py
"""Dispatcher de extracción de texto por tipo de archivo + captura de cite."""
from __future__ import annotations

import base64
import io
import json
import logging
import re
from typing import Optional

import azure.functions as func
import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from .mcp_helpers import ToolProperty, get_arguments, tool_properties_json

logger = logging.getLogger(__name__)

_TOOL_PROPERTIES = tool_properties_json(
    [
        ToolProperty("fileBase64", "string", "Contenido del documento de respaldo en base64.", isRequired=True),
        ToolProperty("fileName", "string", "Nombre del archivo.", isRequired=True),
    ]
)

_CITE_RE = re.compile(r"BANCO\s+FIE\s+S\.?A\.?[\/\s][\w\-\.\/]{3,40}\/\d{4}", re.IGNORECASE)


def register_extract_backup_content_tool(app: func.FunctionApp):
    @app.mcp_tool_trigger(
        arg_name="context",
        tool_name="ExtractBackupContent",
        description="Extrae texto y cite de un documento de respaldo (pdf, docx, pptx, xlsx).",
        tool_properties=_TOOL_PROPERTIES,
    )
    def extract_backup_content(context) -> str:
        try:
            args = get_arguments(context)
            file_b64 = args.get("fileBase64")
            file_name = args.get("fileName", "")
            if not file_b64 or not file_name:
                raise ValueError("Debe enviar 'fileBase64' y 'fileName'.")

            content = base64.b64decode(file_b64)
            extension = file_name.lower().rsplit(".", 1)[-1]

            if extension == "pdf":
                text = _extract_pdf(content)
            elif extension == "docx":
                text = _extract_docx(content)
            elif extension == "pptx":
                text = _extract_pptx(content)
            elif extension in ("xlsx", "xlsm"):
                text = _extract_xlsx(content)
            else:
                raise ValueError(f"Tipo de archivo no soportado: {extension}")

            cite = _find_cite(text)

            return json.dumps(
                {
                    "status": "ok",
                    "fileName": file_name,
                    "tipoArchivo": extension,
                    "cite": cite,
                    "textoExtraido": text[:12000],
                },
                ensure_ascii=False,
            )
        except Exception as exc:
            logger.exception("ExtractBackupContent error")
            return json.dumps({"status": "error", "message": str(exc)}, ensure_ascii=False)

    return extract_backup_content


def _extract_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append(f"[Notas] {slide.notes_slide.notes_text_frame.text}")
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), data_only=True)
    parts = []
    for sheet in wb.worksheets[:3]:
        parts.append(f"[Hoja: {sheet.title}]")
        for row in sheet.iter_rows(max_row=60, values_only=True):
            values = [str(v) for v in row if v is not None]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _find_cite(text: str) -> Optional[str]:
    match = _CITE_RE.search(text)
    return match.group(0).strip() if match else None

# tools/extract_backup_content.py — reemplazar la función _extract_pdf

from .ocr_helper import extract_text_with_ocr_fallback


def _extract_pdf(content: bytes) -> str:
    return extract_text_with_ocr_fallback(content)